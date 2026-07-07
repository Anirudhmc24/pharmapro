import os
from pptx import Presentation
from pptx.util import Pt

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullet_slide(prs, heading, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = heading
    tf = slide.shapes.placeholders[1].text_frame
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(18)
        p.level = 0

def build_deck():
    prs = Presentation()
    add_title_slide(prs,
        "PharmaPro – Smart Pharmacy Management",
        "All‑in‑one solution for inventory, billing, and compliance")
    benefits = [
        "Rapid invoice scanning – multi‑image upload, auto‑extract batch & expiry",
        "One‑click WhatsApp sharing – send bills instantly & store customer contacts",
        "Smart stock management – FEFO, low‑stock & expiry alerts",
        "Master‑database auto‑fill – >250k medicines with composition & MRP",
        "Dynamic shop layout mapping – shelves, boxes, fixtures",
        "Regulatory compliance – GSTR‑1, drug‑interaction warnings, batch traceability",
        "Portable executable – single‑file Windows .exe, no Python install needed"
    ]
    add_bullet_slide(prs, "Why Choose PharmaPro?", benefits)
    screenshots = [
        "Inventory Management – sleek modern UI (see screenshot)",
        "Billing POS – fast search, auto‑batch selection",
        "WhatsApp Send – one‑click bill sharing"
    ]
    add_bullet_slide(prs, "Key Screens", screenshots)
    steps = [
        "Download `PharmaPro.exe` from the `dist/` folder",
        "Run the installer – no admin rights required",
        "Follow the Initial Setup steps in the User Manual",
        "Start managing inventory, billing, and reports instantly"
    ]
    add_bullet_slide(prs, "Getting Started", steps)
    contact = [
        "Website: https://pharmapro.example.com",
        "Email: support@pharmapro.example.com",
        "Phone: +91‑12345‑67890"
    ]
    add_bullet_slide(prs, "Contact Us", contact)
    out_path = os.path.abspath(os.path.join('docs', 'PharmaPro_SlideDeck.pptx'))
    prs.save(out_path)
    print(f'Slide deck generated at {out_path}')

if __name__ == '__main__':
    build_deck()
